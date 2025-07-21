# src/veritas/engine/reporting.py

import io
import pandas as pd
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go
from typing import Dict, Any

# --- PDF Generation Engine ---

class VeritasPDF(FPDF):
    """
    A custom FPDF class with VERITAS-branded headers, footers, and watermarking.
    This provides a standardized, GxP-compliant report format.
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.watermark_text = ""

    def set_watermark(self, text: str) -> None:
        """Sets the text to be used as a watermark on each page."""
        self.watermark_text = text

    def header(self) -> None:
        """Defines the header for each page of the PDF."""
        self.set_font('Helvetica', 'B', 12)
        self.cell(0, 10, 'VERITAS - Automated Data Summary Report', 0, 1, 'C')
        self.set_font('Helvetica', '', 8)
        self.cell(0, 5, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S UTC')}", 0, 1, 'C')
        self.ln(10)

    def footer(self) -> None:
        """Defines the footer for each page of the PDF."""
        self.set_y(-15)
        self.set_font('Helvetica', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
        self.set_x(self.w - self.r_margin - 50) # Position for right-aligned text
        self.cell(50, 10, 'VERITAS CONFIDENTIAL', 0, 0, 'R')

    def _add_watermark(self) -> None:
        """Internal method to render the watermark text on the page."""
        if self.watermark_text:
            self.set_font('Helvetica', 'B', 50)
            self.set_text_color(230, 230, 230) # Light gray
            # Rotate and place the watermark text diagonally in the center
            with self.rotation(45, x=self.w / 2, y=self.h / 2):
                self.text(x=self.w / 2 - 50, y=self.h / 2, text=self.watermark_text)
            self.set_text_color(0, 0, 0) # Reset text color to black

    def add_page(self, orientation: str = '', *args, **kwargs) -> None:
        """Overrides the default add_page to include the watermark automatically."""
        super().add_page(orientation=orientation, *args, **kwargs)
        self._add_watermark()

    def chapter_title(self, title: str) -> None:
        """Adds a formatted chapter title to the PDF."""
        self.set_font('Helvetica', 'B', 12)
        self.set_fill_color(220, 220, 220) # Light gray background
        self.cell(0, 8, title, 0, 1, 'L', fill=True)
        self.ln(4)

    def chapter_body(self, body: str) -> None:
        """Adds a formatted body of text (e.g., commentary) to the PDF."""
        self.set_font('Helvetica', '', 10)
        self.multi_cell(0, 5, body)
        self.ln()

    def add_dataframe(self, df: pd.DataFrame, title: str) -> None:
        """
        Renders a pandas DataFrame as a table in the PDF.
        It uses FPDF2's robust built-in table functionality.
        """
        if not isinstance(df, pd.DataFrame):
            raise TypeError("df must be a pandas DataFrame.")
        if df.empty:
            self.chapter_body(f"No data available for table: '{title}'.")
            return

        self.chapter_title(title)
        self.set_font("Helvetica", size=9)
        
        # Use FPDF2's built-in table functionality for robustness and auto-layout
        with self.table(text_align="CENTER", first_row_as_header=True) as table:
            # Add header from DataFrame columns
            header = table.row()
            for col_name in df.columns:
                header.cell(col_name)
            
            # Add data rows
            for _, data_row in df.iterrows():
                row = table.row()
                for item in data_row:
                    row.cell(str(item))
        self.ln(5)

    def add_plot(self, fig: go.Figure, title: str) -> None:
        """
        Renders a Plotly figure as a high-resolution image in the PDF.

        Args:
            fig (go.Figure): The Plotly figure object to render.
            title (str): A title for the plot section.
        """
        self.chapter_title(title)
        img_bytes = fig.to_image(format="png", width=800, height=450, scale=2)
        img_stream = io.BytesIO(img_bytes)
        
        # Center the image on the page
        page_width = self.w - self.l_margin - self.r_margin
        self.image(img_stream, w=page_width)
        self.ln(5)

    def add_signature_section(self, signature_details: Dict) -> None:
        """Adds a formatted 21 CFR Part 11 compliant electronic signature block."""
        self.chapter_title("Electronic Signature (21 CFR Part 11)")
        sig_body = (
            f"This document was electronically signed and locked in the VERITAS system.\n\n"
            f"Signed By: {signature_details.get('user', 'N/A')}\n"
            f"Signature Timestamp: {signature_details.get('timestamp', 'N/A')}\n"
            f"Meaning of Signature: {signature_details.get('reason', 'N/A')}"
        )
        self.set_font('Helvetica', '', 10)
        self.multi_cell(0, 5, sig_body)
        self.ln()

def generate_pdf_report(report_data: Dict[str, Any], watermark: str = "") -> bytes:
    """
    Top-level function to orchestrate the creation of a formatted PDF report.

    Args:
        report_data (Dict[str, Any]): A dictionary containing all data and configuration for the report.
        watermark (str, optional): Text to display as a watermark (e.g., "DRAFT"). Defaults to "".

    Returns:
        bytes: The generated PDF content as a byte string.
    """
    # Unpack and validate required data
    study_id = report_data.get('study_id')
    commentary = report_data.get('commentary')
    df = report_data.get('data')
    fig = report_data.get('plot_fig')
    sections_config = report_data.get('sections_config')
    signature_details = report_data.get('signature_details')

    if not all([study_id, commentary, isinstance(df, pd.DataFrame), sections_config]):
        raise ValueError("Report data is missing one or more required keys.")

    pdf = VeritasPDF()
    if watermark:
        pdf.set_watermark(watermark)
    pdf.add_page()

    pdf.chapter_title(f"1.0 Summary for Study: {study_id}")
    pdf.chapter_body(f"Analyst Commentary:\n{commentary}")

    if sections_config.get('include_summary_stats'):
        summary_stats = df[report_data['cqa']].describe().round(3).reset_index()
        summary_stats.columns = ['Statistic', 'Value']
        pdf.add_dataframe(summary_stats, "2.1 Summary Statistics")

    if fig and sections_config.get('include_cpk_analysis'):
        plot_title = fig.layout.title.text or "Process Capability Analysis"
        pdf.add_plot(fig, f"2.2 {plot_title}")
    
    if sections_config.get('include_full_dataset'):
        pdf.add_page() # Add a new page for the appendix
        pdf.add_dataframe(df, "3.0 Appendix: Full Dataset")

    if signature_details:
        pdf.add_signature_section(signature_details)

    return pdf.output()

# --- PowerPoint Generation Engine ---

def _add_table_to_slide(slide, df: pd.DataFrame, left: Inches, top: Inches, width: Inches, height: Inches):
    """Helper function to render a pandas DataFrame as a table on a PowerPoint slide."""
    if df.empty:
        return
    rows, cols = df.shape
    rows += 1  # Add a row for the header
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(cols):
        table.columns[i].width = Inches(width.inches / cols)

    for c, col_name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    for r, data_row in enumerate(df.itertuples(index=False)):
        for c, val in enumerate(data_row):
            cell = table.cell(r + 1, c)
            cell.text = str(val) if pd.notna(val) else ""
            cell.text_frame.paragraphs[0].font.size = Pt(11)

def generate_ppt_report(report_data: Dict[str, Any]) -> bytes:
    """
    Generates a standard PowerPoint report with data and plots.

    Args:
        report_data (Dict[str, Any]): A dictionary containing data and configuration for the PPT.

    Returns:
        bytes: The generated .pptx file content as a byte string.
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "VERITAS Automated Study Report"
    slide.placeholders[1].text = (
        f"Study ID: {report_data.get('study_id', 'N/A')}\n"
        f"Generated on: {datetime.now().strftime('%Y-%m-%d')}"
    )

    # Slide 2: Summary Statistics
    if report_data.get('sections_config', {}).get('include_summary_stats'):
        content_slide_layout = prs.slide_layouts[5] # Title and Content layout
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Summary Statistics"
        df = report_data.get('data')
        cqa = report_data.get('cqa')
        if df is not None and cqa:
            summary_stats = df[cqa].describe().round(3).reset_index()
            summary_stats.columns = ['Statistic', 'Value']
            _add_table_to_slide(slide, summary_stats, Inches(1), Inches(1.5), Inches(8), Inches(4))

    # Slide 3: Primary Plot
    fig = report_data.get('plot_fig')
    if fig and report_data.get('sections_config', {}).get('include_cpk_analysis'):
        content_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = fig.layout.title.text or "Analysis Chart"
        img_bytes = fig.to_image(format="png", width=800, height=450, scale=2)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))

    # Save presentation to a byte stream to be returned
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue()
